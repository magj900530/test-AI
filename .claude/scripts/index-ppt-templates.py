"""PPT 模版索引生成器

扫描 D:/AI项目/PPT模版/ 目录，为每个 .pptx/.potx 生成缩略图和元数据，
输出 INDEX.md 索引文件。幂等全量重建——每次运行完整刷新。

Usage:
    python index-ppt-templates.py
"""

import os
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

# --- 配置 ---
TEMPLATE_ROOT = Path("D:/AI项目/PPT模版")
THUMBNAIL_DIR_NAME = "thumbnails"
INDEX_FILENAME = "INDEX.md"
THUMBNAIL_COLS = 2          # 缩略图网格列数（2×2=4页预览）
THUMBNAIL_ROWS = 2
THUMBNAIL_WIDTH = 300       # 每格宽度
CONVERSION_DPI = 100        # PDF 转图片 DPI
JPEG_QUALITY = 85
GRID_PADDING = 10
BORDER_COLOR = "#CCCCCC"

# --- 依赖检查 ---
HAS_PPTX = False
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    pass

HAS_MARKITDOWN = False
try:
    from markitdown import MarkItDown
    HAS_MARKITDOWN = True
except ImportError:
    pass

HAS_PIL = False
try:
    from PIL import Image, ImageDraw, ImageFont
    HAS_PIL = True
except ImportError:
    pass


def _check_tool(name):
    """检查命令行工具是否可用"""
    import shutil
    return shutil.which(name) is not None


# PowerPoint COM（需要 Office 2010+ 才支持 Slide.Export）
HAS_POWERPOINT = False
if sys.platform == "win32":
    try:
        import win32com.client
        ppt_test = win32com.client.Dispatch("PowerPoint.Application")
        ver = float(ppt_test.Version)
        ppt_test.Quit()
        # Slide.Export 方法从 Office 2010 (v14.0) 开始支持
        if ver >= 14.0:
            HAS_POWERPOINT = True
    except Exception:
        pass

HAS_SOFFICE = _check_tool("soffice")
HAS_PDFTOPPM = _check_tool("pdftoppm")
HAS_THUMBNAIL_TOOLS = HAS_POWERPOINT or (HAS_SOFFICE and HAS_PDFTOPPM and HAS_PIL) or (HAS_PPTX and HAS_PIL)


def find_soffice_env():
    """获取 LibreOffice 环境变量（跨平台安全）"""
    import platform
    env = os.environ.copy()
    env["SAL_USE_VCLPLUGIN"] = "svp"

    # Windows 不需要 Unix socket shim
    if platform.system() == "Windows":
        return env

    # Linux: 尝试加载 soffice.py 的 shim（处理沙箱 AF_UNIX 限制）
    soffice_py = Path.home() / ".claude/skills/pptx/scripts/office/soffice.py"
    if soffice_py.exists():
        try:
            import importlib.util
            spec = importlib.util.spec_from_file_location("soffice", str(soffice_py))
            soffice = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(soffice)
            if hasattr(soffice, "get_soffice_env"):
                return soffice.get_soffice_env()
        except Exception:
            pass  # 加载失败就返回基础 env

    return env


def scan_templates(root):
    """递归扫描，返回 [(相对分类, 文件路径, 文件名), ...]"""
    templates = []
    if not root.exists():
        print(f"[ERROR] Template directory not found: {root}")
        return templates

    for f in sorted(root.rglob("*")):
        if f.is_file() and f.suffix.lower() in (".pptx", ".potx"):
            category = f.parent.relative_to(root).as_posix()
            if category == ".":
                category = "未分类"
            templates.append((category, f, f.name))
    return templates


def get_slide_count(pptx_path):
    """用 zipfile 直接读 XML 获取页数（不依赖 python-pptx）"""
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            pres_xml = zf.read("ppt/presentation.xml").decode("utf-8")
            return pres_xml.count('<p:sldId ')
    except Exception:
        return "?"


def get_text_preview(pptx_path, max_chars=300):
    """提取模板文本摘要"""
    text = ""

    # 优先用 markitdown（提取效果最好）
    if HAS_MARKITDOWN:
        try:
            md = MarkItDown()
            result = md.convert(str(pptx_path))
            text = result.text_content[:max_chars]
        except Exception:
            pass

    # 其次用 python-pptx
    if not text and HAS_PPTX:
        try:
            prs = Presentation(str(pptx_path))
            lines = []
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t and len(t) > 2:
                                lines.append(t)
            text = "\n".join(lines)[:max_chars]
        except Exception:
            pass

    # 最后用 zipfile 读 slide XML（兜底）
    if not text:
        try:
            with zipfile.ZipFile(pptx_path, "r") as zf:
                slides = sorted(
                    [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")],
                    key=lambda x: int(x.replace("ppt/slides/slide", "").replace(".xml", ""))
                )
                lines = []
                for slide_path in slides[:4]:  # 只看前 4 页
                    content = zf.read(slide_path).decode("utf-8", errors="replace")
                    # 提取所有 <a:t> 文本
                    import re
                    texts = re.findall(r'<a:t[^>]*>(.*?)</a:t>', content)
                    for t in texts:
                        t = t.strip()
                        if t and len(t) > 2:
                            lines.append(t)
                text = "\n".join(lines)[:max_chars]
        except Exception:
            text = "(无法提取文本)"

    return text.strip() or "(空白模板)"


def generate_thumbnail(pptx_path, output_path):
    """生成缩略图网格（2×2）

    优先使用 PowerPoint COM（Windows），后备 LibreOffice（跨平台）。
    """
    if not HAS_THUMBNAIL_TOOLS:
        return False

    # 按优先级尝试：PowerPoint > LibreOffice > python-pptx
    if HAS_POWERPOINT:
        ok = _thumb_via_powerpoint(pptx_path, output_path)
        if ok:
            return True
    if HAS_SOFFICE and HAS_PDFTOPPM:
        ok = _thumb_via_libreoffice(pptx_path, output_path)
        if ok:
            return True
    if HAS_PPTX and HAS_PIL:
        return _thumb_via_python(pptx_path, output_path)
    return False


def _thumb_via_powerpoint(pptx_path, output_path):
    """Windows: 用 PowerPoint COM 导出幻灯片为图片"""
    import pythoncom
    import win32com.client
    import shutil

    pythoncom.CoInitialize()
    ppt = None
    presentation = None
    tmp_copy = None
    try:
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        ppt.Visible = False

        # PowerPoint 2007 对非 ASCII 路径支持差，复制到临时英文路径
        use_path = pptx_path
        try:
            pptx_path.name.encode("ascii")
        except UnicodeEncodeError:
            tmp_copy = Path(tempfile.gettempdir()) / ("_ppt_thumb_" + str(hash(str(pptx_path))) + pptx_path.suffix)
            shutil.copy2(pptx_path, tmp_copy)
            use_path = tmp_copy

        presentation = ppt.Presentations.Open(str(use_path), WithWindow=False)

        with tempfile.TemporaryDirectory() as tmp:
            tmp = Path(tmp)
            image_paths = []
            max_slides = min(presentation.Slides.Count, THUMBNAIL_COLS * THUMBNAIL_ROWS)
            for i in range(1, max_slides + 1):
                img_path = tmp / f"slide-{i:02d}.png"
                try:
                    # Slide.Export(FileName, FilterName, ScaleWidth, ScaleHeight)
                    presentation.Slides(i).Export(str(img_path), "PNG", 1920, 1080)
                except Exception:
                    # 旧版 PowerPoint 回退：尝试不指定尺寸
                    presentation.Slides(i).Export(str(img_path), "PNG")
                if img_path.exists():
                    image_paths.append(img_path)

            presentation.Close()
            presentation = None
            ppt.Quit()
            ppt = None

            # 清理临时文件
            if tmp_copy and tmp_copy.exists():
                tmp_copy.unlink()

            if not image_paths:
                return False
            return _make_grid(image_paths, output_path)

    except Exception as e:
        print(f"  [WARN] PowerPoint export failed: {e}")
        return False
    finally:
        try:
            if presentation:
                presentation.Close()
        except Exception:
            pass
        try:
            if ppt:
                ppt.Quit()
        except Exception:
            pass
        try:
            if tmp_copy and tmp_copy.exists():
                tmp_copy.unlink()
        except Exception:
            pass
        pythoncom.CoUninitialize()


def _thumb_via_libreoffice(pptx_path, output_path):
    """Linux/跨平台: LibreOffice + pdftoppm + Pillow"""
    if not (HAS_SOFFICE and HAS_PDFTOPPM and HAS_PIL):
        return False

    env = find_soffice_env()
    with tempfile.TemporaryDirectory() as tmp:
        tmp = Path(tmp)
        pdf_path = tmp / "temp.pdf"

        result = subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(tmp), str(pptx_path)],
            capture_output=True, text=True, env=env,
        )
        if result.returncode != 0 or not pdf_path.exists():
            print(f"  [WARN] PDF conversion failed: {pptx_path.name}")
            return False

        result = subprocess.run(
            ["pdftoppm", "-jpeg", "-r", str(CONVERSION_DPI), str(pdf_path), str(tmp / "slide")],
            capture_output=True, text=True,
        )
        images = sorted(tmp.glob("slide-*.jpg"))
        if not images:
            print(f"  [WARN] Image conversion failed: {pptx_path.name}")
            return False

        return _make_grid(images[:THUMBNAIL_COLS * THUMBNAIL_ROWS], output_path)


def _thumb_via_python(pptx_path, output_path):
    """纯 Python 缩略图：用 python-pptx + Pillow 提取配色和版式生成预览

    不依赖任何外部渲染引擎，展示：
    - 幻灯片背景色和主题色
    - 每页标题文字和位置
    - 形状/图片的简化轮廓
    """
    if not HAS_PPTX or not HAS_PIL:
        return False

    try:
        tmp_renamed = None
        # .potx 不被 python-pptx 识别（MIME type 不同），修补 [Content_Types].xml
        load_path = pptx_path
        if pptx_path.suffix.lower() == ".potx":
            import shutil, io
            tmp_renamed = Path(tempfile.gettempdir()) / ("_thumb_" + pptx_path.stem + ".pptx")
            # 读原始 zip，替换 MIME type，写为新文件
            with zipfile.ZipFile(pptx_path, "r") as zin:
                with zipfile.ZipFile(tmp_renamed, "w", zipfile.ZIP_DEFLATED) as zout:
                    for item in zin.infolist():
                        data = zin.read(item.filename)
                        if item.filename == "[Content_Types].xml":
                            data = data.replace(
                                b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                                b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
                            )
                        zout.writestr(item, data)
            load_path = tmp_renamed

        prs = Presentation(str(load_path))
        # 读取幻灯片尺寸（EMU → px @ 96 DPI）
        sw_emu = prs.slide_width   # EMU
        sh_emu = prs.slide_height
        sw_px = int(sw_emu / 9525)  # 1 px = 9525 EMU at 96 DPI
        sh_px = int(sh_emu / 9525)

        previews = []
        max_slides = min(len(prs.slides), THUMBNAIL_COLS * THUMBNAIL_ROWS)

        for idx in range(max_slides):
            slide = prs.slides[idx]
            img = _render_slide_preview(slide, sw_px, sh_px)
            previews.append(img)

        if not previews:
            return False

        # 拼 2×2 网格
        result = _make_pil_grid(previews, output_path)

        # 清理临时文件
        if tmp_renamed and tmp_renamed.exists():
            tmp_renamed.unlink()

        return result

    except Exception as e:
        print(f"  [WARN] Python thumbnail failed: {e}")
        if tmp_renamed and tmp_renamed.exists():
            tmp_renamed.unlink()
        return False


def _render_slide_preview(slide, sw_px, sh_px):
    """用 Pillow 渲染单页幻灯片预览"""
    # 缩放到缩略图宽度
    scale = THUMBNAIL_WIDTH / sw_px
    w, h = THUMBNAIL_WIDTH, int(sh_px * scale)

    img = Image.new("RGB", (w, h), "#FFFFFF")
    draw = ImageDraw.Draw(img)

    # 尝试获取背景色
    bg_color = _get_slide_bg(slide)

    # 绘制背景
    draw.rectangle([0, 0, w, h], fill=bg_color)

    # 遍历形状，绘制简化版
    for shape in slide.shapes:
        sx = int(shape.left / 9525 * scale)
        sy = int(shape.top / 9525 * scale)
        sw = int(shape.width / 9525 * scale)
        sh_val = int(shape.height / 9525 * scale)

        fill_color = None
        if hasattr(shape, "fill"):
            try:
                fill_color = _get_fill_color(shape.fill, bg_color)
            except Exception:
                pass

        if shape.has_text_frame:
            # 文本框：浅灰底 + 文字
            if fill_color:
                draw.rectangle([sx, sy, sx + sw, sy + sh_val], fill=fill_color, outline="#DDDDDD")
            first_text = ""
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    first_text = t[:40]
                    break
            if first_text:
                font_size = max(6, int(10 * scale))
                try:
                    font = ImageFont.truetype("arial.ttf", size=font_size)
                except Exception:
                    font = ImageFont.load_default()
                # 截断避免溢出
                draw.text((sx + 2, sy + 2), first_text, fill="#333333", font=font)
        elif shape.shape_type == 13:  # Picture
            draw.rectangle([sx, sy, sx + sw, sy + sh_val], fill="#E8ECF0", outline="#CCCCCC")
            draw.line([sx, sy, sx + sw, sy + sh_val], fill="#CCCCCC")
            draw.line([sx + sw, sy, sx, sy + sh_val], fill="#CCCCCC")
        elif fill_color and fill_color != bg_color:
            draw.rectangle([sx, sy, sx + sw, sy + sh_val], fill=fill_color, outline="#DDDDDD")

    return img


def _get_slide_bg(slide):
    """提取幻灯片背景色"""
    try:
        bg = slide.background
        if bg.fill.type is not None:
            return _get_fill_color(bg.fill, "#FFFFFF")
    except Exception:
        pass

    # 从 slide layout 继承
    try:
        layout_bg = slide.slide_layout.background
        if layout_bg.fill.type is not None:
            return _get_fill_color(layout_bg.fill, "#FFFFFF")
    except Exception:
        pass

    return "#FFFFFF"


def _get_fill_color(fill, default="#FFFFFF"):
    """从 pptx Fill 对象提取颜色"""
    try:
        if fill.type == 1:  # Solid
            fc = fill.fore_color
            if fc.type == 1:  # RGB
                return f"#{fc.rgb:06X}" if fc.rgb else default
            elif fc.type == 2:  # Theme
                return f"#{fc.theme_color:06X}" if fc.theme_color else default
    except Exception:
        pass
    return default


def _make_pil_grid(pil_images, output_path):
    """拼接 Pillow Image 列表为网格并保存"""
    w, h = pil_images[0].size
    cols = THUMBNAIL_COLS
    rows = min((len(pil_images) + cols - 1) // cols, THUMBNAIL_ROWS)
    actual = min(len(pil_images), cols * rows)

    grid_w = cols * w + (cols + 1) * GRID_PADDING
    grid_h = rows * h + (rows + 1) * GRID_PADDING

    grid = Image.new("RGB", (grid_w, grid_h), "white")
    for i in range(actual):
        row, col = i // cols, i % cols
        x = col * w + (col + 1) * GRID_PADDING
        y = row * h + (row + 1) * GRID_PADDING
        grid.paste(pil_images[i], (x, y))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    grid.save(str(output_path), quality=JPEG_QUALITY)
    return True


def _make_grid(image_paths, output_path):
    """拼接缩略图网格"""
    try:
        with Image.open(image_paths[0]) as img:
            aspect = img.height / img.width
        cell_h = int(THUMBNAIL_WIDTH * aspect)

        cols = THUMBNAIL_COLS
        rows = min((len(image_paths) + cols - 1) // cols, THUMBNAIL_ROWS)
        actual = min(len(image_paths), cols * rows)

        grid_w = cols * THUMBNAIL_WIDTH + (cols + 1) * GRID_PADDING
        grid_h = rows * cell_h + (rows + 1) * GRID_PADDING

        grid = Image.new("RGB", (grid_w, grid_h), "white")
        draw = ImageDraw.Draw(grid)

        for i in range(actual):
            row, col = i // cols, i % cols
            x = col * THUMBNAIL_WIDTH + (col + 1) * GRID_PADDING
            y = row * cell_h + (row + 1) * GRID_PADDING

            with Image.open(image_paths[i]) as img:
                img.thumbnail((THUMBNAIL_WIDTH, cell_h), Image.Resampling.LANCZOS)
                w, h = img.size
                tx = x + (THUMBNAIL_WIDTH - w) // 2
                ty = y + (cell_h - h) // 2
                grid.paste(img, (tx, ty))
                draw.rectangle(
                    [tx - 1, ty - 1, tx + w, ty + h],
                    outline=BORDER_COLOR, width=1,
                )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        grid.save(str(output_path), quality=JPEG_QUALITY)
        return True
    except Exception as e:
        print(f"  [WARN] Thumbnail generation failed: {e}")
        return False


def format_size(size_bytes):
    """格式化文件大小"""
    for unit in ["B", "KB", "MB"]:
        if size_bytes < 1024:
            return f"{size_bytes:.0f} {unit}"
        size_bytes /= 1024
    return f"{size_bytes:.1f} GB"


def generate_index(templates, thumb_dir):
    """生成 INDEX.md"""
    lines = [
        "# PPT 模版索引",
        "",
        f"> 自动生成于 {_now()}，共 {len(templates)} 个模版。",
        f"> 增删模版后运行 `python {__file__}` 刷新。",
        "",
    ]
    if not HAS_THUMBNAIL_TOOLS:
        lines.append("> ⚠️ 缩略图功能不可用（需安装 python-pptx + Pillow），索引仅含文字信息。")
    lines += ["", "---", ""]

    # 按分类分组
    by_category = {}
    for cat, path, name in templates:
        by_category.setdefault(cat, []).append((path, name))

    for cat in sorted(by_category.keys()):
        items = by_category[cat]
        lines.append(f"## {cat}（{len(items)} 个）")
        lines.append("")

        for path, name in items:
            slide_count = get_slide_count(path)
            file_size = format_size(path.stat().st_size)
            fmt = "POTX 模版" if path.suffix.lower() == ".potx" else "PPTX"

            # 缩略图（如可用）
            thumb_name = path.stem + ".jpg"
            thumb_rel = f"{THUMBNAIL_DIR_NAME}/{thumb_name}"
            thumb_file = thumb_dir / thumb_name

            lines.append(f"### {path.stem}")
            lines.append("")
            lines.append(f"| 属性 | 值 |")
            lines.append(f"|------|-----|")
            lines.append(f"| 文件名 | `{name}` |")
            lines.append(f"| 格式 | {fmt} |")
            lines.append(f"| 大小 | {file_size} |")
            lines.append(f"| 页数 | {slide_count} 页 |")
            if thumb_file.exists():
                lines.append(f"| 预览 | ![{path.stem}]({thumb_rel}) |")
            lines.append("")

            # 文本摘要
            preview = get_text_preview(path)
            lines.append(f"**内容摘要：**")
            lines.append("")
            for line in preview.split("\n")[:6]:
                if line.strip():
                    lines.append(f"> {line.strip()}")
            lines.append("")
            lines.append("---")
            lines.append("")

    return "\n".join(lines)


def _now():
    from datetime import datetime
    return datetime.now().strftime("%Y-%m-%d %H:%M")


def main():
    print(f"Scanning: {TEMPLATE_ROOT}")
    templates = scan_templates(TEMPLATE_ROOT)

    if not templates:
        print("No templates found, generating empty index.")
        # 仍生成一个空索引
        index_path = TEMPLATE_ROOT / INDEX_FILENAME
        index_path.write_text(
            "# PPT 模版索引\n\n> 暂无模版。请将 .pptx/.potx 文件放入此目录的子文件夹中。\n",
            encoding="utf-8"
        )
        return

    print(f"Found {len(templates)} template files")
    if HAS_POWERPOINT:
        print("(thumbnails: PowerPoint COM)")
    elif HAS_SOFFICE and HAS_PDFTOPPM:
        print("(thumbnails: LibreOffice)")
    elif HAS_PPTX and HAS_PIL:
        print("(thumbnails: python-pptx)")
    print()

    thumb_dir = TEMPLATE_ROOT / THUMBNAIL_DIR_NAME
    thumb_dir.mkdir(parents=True, exist_ok=True)

    # 清理孤儿缩略图
    valid_thumbs = {t[1].stem + ".jpg" for t in templates}
    for old_thumb in thumb_dir.glob("*.jpg"):
        if old_thumb.name not in valid_thumbs:
            old_thumb.unlink()
            print(f"Cleanup orphan thumb: {old_thumb.name}")

    # 处理每个模版
    for i, (cat, path, name) in enumerate(templates, 1):
        thumb_path = thumb_dir / (path.stem + ".jpg")
        status = ""

        # 缩略图
        if HAS_THUMBNAIL_TOOLS:
            if not thumb_path.exists():
                print(f"[{i}/{len(templates)}] {cat}/{name} -> generating thumb...")
                if generate_thumbnail(path, thumb_path):
                    status = "OK"
                else:
                    status = "ERR(thumb)"
            else:
                print(f"[{i}/{len(templates)}] {cat}/{name} -> skip (thumb exists)")
                status = "OK"
        else:
            print(f"[{i}/{len(templates)}] {cat}/{name}")
            status = "OK"

        if status:
            print(f"  [{status}] {get_slide_count(path)} slides, {format_size(path.stat().st_size)}")

    # 生成索引
    print(f"\nGenerating INDEX.md...")
    index_content = generate_index(templates, thumb_dir)
    index_path = TEMPLATE_ROOT / INDEX_FILENAME
    index_path.write_text(index_content, encoding="utf-8")
    print(f"Done -> {index_path}")


if __name__ == "__main__":
    main()
